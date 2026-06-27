"""Modern threaded comments for PowerPoint (.pptx) — list / add / reply / resolve / delete.

Confirmed against real PowerPoint (build 16.0). Modern comments use the p188 schema:

    ppt/authors.xml                         <p188:author id name initials userId providerId>
    ppt/comments/modernComment{n}.xml       <p188:cmLst> of <p188:cm> per slide
    presentation rel  -> authors.xml
    slide rel         -> the slide's modern comment part

A comment anchors to its slide via <pc:sldMk cId=.. sldId=..>: cId = the slide's creationId
(in the slide extLst), sldId = the slide's id in presentation.xml. Threading is NESTED — a
reply is a <p188:reply> inside the comment's <p188:replyLst>. Resolve = status="resolved" on
the <p188:cm>.
"""
from __future__ import annotations

import random
import re
import zipfile

from lxml import etree

from _ooxml_zip import patch_parts, _rels_name
from _errors import AnchorNotFound, CommentNotFound
from _util import guid, utc_now, initials as _initials_of

# --- namespaces ---
P188 = "http://schemas.microsoft.com/office/powerpoint/2018/8/main"
PC = "http://schemas.microsoft.com/office/powerpoint/2013/main/command"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"

CT_COMMENTS = "application/vnd.ms-powerpoint.comments+xml"
CT_AUTHORS = "application/vnd.ms-powerpoint.authors+xml"
REL_COMMENTS = "http://schemas.microsoft.com/office/2018/10/relationships/comments"
REL_AUTHORS = "http://schemas.microsoft.com/office/2018/10/relationships/authors"

AUTHORS_PART = "ppt/authors.xml"
CREATIONID_URI = "{BB962C8B-B14F-4D97-AF65-F5344CB8AC3E}"


def _ppt_dt(date):
    if date:
        return date
    n = utc_now()
    return n.strftime("%Y-%m-%dT%H:%M:%S.") + f"{n.microsecond // 1000:03d}"


# ---------------------------------------------------------------------------
# OPC relationship helpers (driven by a get_bytes(name)->bytes|None callable)
# ---------------------------------------------------------------------------

def _rels_by_type(get_bytes, part):
    raw = get_bytes(_rels_name(part))
    by_id, by_type = {}, {}
    if raw:
        for rel in etree.fromstring(raw):
            by_id[rel.get("Id")] = rel.get("Target")
            by_type.setdefault(rel.get("Type"), []).append((rel.get("Id"), rel.get("Target")))
    return by_id, by_type


def _abs_target(base_part, target):
    if target.startswith("/"):
        return target[1:]
    base = base_part.rsplit("/", 1)[0]
    out = []
    for p in (base + "/" + target).split("/"):
        if p == "..":
            if out:
                out.pop()
        elif p not in ("", "."):
            out.append(p)
    return "/".join(out)


def _rel_target(get_bytes, part, rel_type):
    _, by_type = _rels_by_type(get_bytes, part)
    rels = by_type.get(rel_type)
    return _abs_target(part, rels[0][1]) if rels else None


def _slide_number(slide_part):
    m = re.search(r"slide(\d+)\.xml$", slide_part)
    return m.group(1) if m else "1"


def _slides(get_bytes):
    """Ordered list of (sldId, slide_part) for the presentation."""
    pres = etree.fromstring(get_bytes("ppt/presentation.xml"))
    by_id, _ = _rels_by_type(get_bytes, "ppt/presentation.xml")
    out = []
    lst = pres.find(f"{{{P}}}sldIdLst")
    if lst is not None:
        for sid in lst.findall(f"{{{P}}}sldId"):
            tgt = by_id.get(sid.get(f"{{{R}}}id"))
            if tgt:
                out.append((sid.get("id"), _abs_target("ppt/presentation.xml", tgt)))
    return out


# ---------------------------------------------------------------------------
# Reading
# ---------------------------------------------------------------------------

def _authors_map(get_bytes):
    raw = get_bytes(AUTHORS_PART)
    out = {}
    if raw:
        for a in etree.fromstring(raw).iter(f"{{{P188}}}author"):
            out[a.get("id")] = a.get("name")
    return out


def list_comments(path) -> list[dict]:
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())
        get = lambda n: (z.read(n) if n in names else None)  # noqa: E731
        authors = _authors_map(get)
        titles = {}
        try:
            from pptx import Presentation
            prs = Presentation(path)
            for i, sl in enumerate(prs.slides, 1):
                try:
                    titles[i] = sl.shapes.title.text if sl.shapes.title else None
                except Exception:
                    titles[i] = None
        except Exception:
            titles = {}

        records = []
        for idx, (sld_id, slide_part) in enumerate(_slides(get), 1):
            cpart = _rel_target(get, slide_part, REL_COMMENTS)
            raw = get(cpart) if cpart else None
            if not raw:
                continue
            root = etree.fromstring(raw)
            for cm in root.findall(f"{{{P188}}}cm"):
                cid = cm.get("id")
                resolved = cm.get("status") == "resolved"
                ctx = titles.get(idx) or ""
                records.append({
                    "id": cid, "thread_id": cid, "parent_id": None, "is_reply": False,
                    "author": authors.get(cm.get("authorId")), "date": cm.get("created"),
                    "text": _own_text(cm), "resolved": resolved,
                    "anchor": {"kind": "pptx", "slide": idx}, "anchor_text": ctx,
                    "context": (f"Slide {idx}: {ctx}" if ctx else f"Slide {idx}"),
                    "location": f"Slide {idx}",
                })
                rl = cm.find(f"{{{P188}}}replyLst")
                if rl is not None:
                    for rep in rl.findall(f"{{{P188}}}reply"):
                        records.append({
                            "id": rep.get("id"), "thread_id": cid, "parent_id": cid, "is_reply": True,
                            "author": authors.get(rep.get("authorId")), "date": rep.get("created"),
                            "text": _own_text(rep), "resolved": resolved,
                            "anchor": {"kind": "pptx", "slide": idx}, "anchor_text": ctx,
                            "context": (f"Slide {idx}: {ctx}" if ctx else f"Slide {idx}"),
                            "location": f"Slide {idx}",
                        })
    return records


def _own_text(node):
    """Text of a cm/reply's OWN txBody (direct child), excluding nested replies."""
    tb = node.find(f"{{{P188}}}txBody")
    return "".join(t.text or "" for t in tb.iter(f"{{{A}}}t")) if tb is not None else ""


# ---------------------------------------------------------------------------
# Writing
# ---------------------------------------------------------------------------

def _all_author_ids(ps):
    if ps.has(AUTHORS_PART):
        return {a.get("id") for a in ps.get_xml(AUTHORS_PART).iter(f"{{{P188}}}author")}
    return set()


def _ensure_author(ps, name, initials):
    if ps.has(AUTHORS_PART):
        root = ps.get_xml(AUTHORS_PART)
        for a in root.findall(f"{{{P188}}}author"):
            if a.get("name") == name:
                return a.get("id")
    else:
        root = etree.Element(f"{{{P188}}}authorLst", nsmap={"a": A, "r": R, "p188": P188})
        ps.add_xml_part(AUTHORS_PART, root)
        ps.ensure_content_type("/" + AUTHORS_PART, CT_AUTHORS)
        ps.rels_for("ppt/presentation.xml").add_rel(REL_AUTHORS, "authors.xml")
    aid = guid(exclude=_all_author_ids(ps))
    a = etree.SubElement(root, f"{{{P188}}}author")
    a.set("id", aid)
    a.set("name", name)
    a.set("initials", initials or _initials_of(name))
    a.set("userId", name)
    a.set("providerId", "None")
    ps.set_xml(AUTHORS_PART, root)
    return aid


def _all_creation_ids(ps):
    """Every slide's existing creationId, so a freshly-minted one stays unique across the deck."""
    used = set()
    for _sid, sp in _slides(ps.get_bytes):
        try:
            ext = ps.get_xml(sp).find(f"{{{P}}}extLst")
        except KeyError:
            continue
        if ext is not None:
            for e in ext.findall(f"{{{P}}}ext"):
                if e.get("uri") == CREATIONID_URI:
                    cid = e.find(f"{{{P14}}}creationId")
                    if cid is not None and cid.get("val"):
                        used.add(cid.get("val"))
    return used


def _slide_creation_id(ps, slide_part):
    """The slide's creationId (cId for the comment anchor); create one if the slide lacks it."""
    root = ps.get_xml(slide_part)
    ext = root.find(f"{{{P}}}extLst")
    if ext is not None:
        for e in ext.findall(f"{{{P}}}ext"):
            if e.get("uri") == CREATIONID_URI:
                cid = e.find(f"{{{P14}}}creationId")
                if cid is not None and cid.get("val"):
                    return cid.get("val")
    used = _all_creation_ids(ps)
    val = str(random.getrandbits(32) or 1)
    while val in used:
        val = str(random.getrandbits(32) or 1)
    if ext is None:
        ext = etree.SubElement(root, f"{{{P}}}extLst")   # extLst is the last child of the slide
    e = etree.SubElement(ext, f"{{{P}}}ext")
    e.set("uri", CREATIONID_URI)
    cid = etree.SubElement(e, f"{{{P14}}}creationId", nsmap={"p14": P14})
    cid.set("val", val)
    ps.set_xml(slide_part, root)
    return val


def _comments_part(ps, slide_part):
    cpart = _rel_target(ps.get_bytes, slide_part, REL_COMMENTS)
    if cpart and ps.has(cpart):
        return cpart, ps.get_xml(cpart)
    if not cpart:
        n = _slide_number(slide_part)
        cpart = f"ppt/comments/modernComment{n}.xml"
        i = int(n)
        while ps.has(cpart):
            i += 1
            cpart = f"ppt/comments/modernComment{i}.xml"
        ps.rels_for(slide_part).add_rel(REL_COMMENTS, "../comments/" + cpart.rsplit("/", 1)[1])
        ps.ensure_content_type("/" + cpart, CT_COMMENTS)
    root = etree.Element(f"{{{P188}}}cmLst", nsmap={"a": A, "r": R, "p188": P188})
    ps.add_xml_part(cpart, root)
    return cpart, root


def _txbody(text):
    tb = etree.Element(f"{{{P188}}}txBody")
    etree.SubElement(tb, f"{{{A}}}bodyPr")
    etree.SubElement(tb, f"{{{A}}}lstStyle")
    p = etree.SubElement(tb, f"{{{A}}}p")
    r = etree.SubElement(p, f"{{{A}}}r")
    etree.SubElement(r, f"{{{A}}}rPr").set("lang", "en-US")
    etree.SubElement(r, f"{{{A}}}t").text = text
    return tb


def _find_cm(ps, cid):
    """Return (slide_idx, slide_part, cpart, cmroot, cm_element) for the thread that `cid`
    belongs to (cid may be the cm id or a reply id), else None."""
    for idx, (sld_id, slide_part) in enumerate(_slides(ps.get_bytes), 1):
        cpart = _rel_target(ps.get_bytes, slide_part, REL_COMMENTS)
        if cpart and ps.has(cpart):
            root = ps.get_xml(cpart)
            for cm in root.findall(f"{{{P188}}}cm"):
                if cm.get("id") == cid or any(r.get("id") == cid for r in cm.iter(f"{{{P188}}}reply")):
                    return idx, slide_part, cpart, root, cm
    return None


def add_comment(path, anchor, text, *, author, initials=None, date=None) -> str:
    """Add a top-level comment to a slide. `anchor` = {"slide": N} (1-based)."""
    out = {}

    def mut(ps):
        slides = _slides(ps.get_bytes)
        n = anchor.get("slide", 1)
        if not (1 <= n <= len(slides)):
            raise AnchorNotFound(f"slide {n} out of range (1..{len(slides)})")
        sld_id, slide_part = slides[n - 1]
        aid = _ensure_author(ps, author, initials)
        cid_val = _slide_creation_id(ps, slide_part)
        cpart, cmroot = _comments_part(ps, slide_part)
        new_id = guid()
        cm = etree.SubElement(cmroot, f"{{{P188}}}cm")
        cm.set("id", new_id)
        cm.set("authorId", aid)
        cm.set("created", _ppt_dt(date))
        smk = etree.SubElement(cm, f"{{{PC}}}sldMkLst", nsmap={"pc": PC})
        etree.SubElement(smk, f"{{{PC}}}docMk")
        sm = etree.SubElement(smk, f"{{{PC}}}sldMk")
        sm.set("cId", cid_val)
        sm.set("sldId", sld_id)
        pos = etree.SubElement(cm, f"{{{P188}}}pos")
        pos.set("x", "1270000")
        pos.set("y", "1270000")
        cm.append(_txbody(text))
        ps.set_xml(cpart, cmroot)
        out["id"] = new_id

    patch_parts(path, mut)
    return out["id"]


def reply(path, parent_id, text, *, author, initials=None, date=None) -> str:
    out = {}

    def mut(ps):
        found = _find_cm(ps, parent_id)
        if found is None:
            raise CommentNotFound(f"no comment with id {parent_id}")
        _, _, cpart, cmroot, cm = found
        aid = _ensure_author(ps, author, initials)
        new_id = guid()
        rep = etree.Element(f"{{{P188}}}reply")
        rep.set("id", new_id)
        rep.set("authorId", aid)
        rep.set("created", _ppt_dt(date))
        rep.append(_txbody(text))
        rl = cm.find(f"{{{P188}}}replyLst")
        if rl is None:
            rl = etree.Element(f"{{{P188}}}replyLst")
            cm.find(f"{{{P188}}}txBody").addprevious(rl)   # replyLst sits before the cm's txBody
        rl.append(rep)
        ps.set_xml(cpart, cmroot)
        out["id"] = new_id

    patch_parts(path, mut)
    return out["id"]


def set_status(path, comment_id, resolved: bool) -> None:
    def mut(ps):
        found = _find_cm(ps, comment_id)
        if found is None:
            raise CommentNotFound(f"no comment with id {comment_id}")
        _, _, cpart, cmroot, cm = found
        if resolved:
            cm.set("status", "resolved")
        elif "status" in cm.attrib:
            del cm.attrib["status"]
        ps.set_xml(cpart, cmroot)

    patch_parts(path, mut)


def delete(path, comment_id) -> None:
    """Delete a comment. A reply goes alone; a thread root takes its whole thread."""
    def mut(ps):
        found = _find_cm(ps, comment_id)
        if found is None:
            raise CommentNotFound(f"no comment with id {comment_id}")
        _, _, cpart, cmroot, cm = found
        if cm.get("id") == comment_id:
            cmroot.remove(cm)                       # whole thread
        else:
            rl = cm.find(f"{{{P188}}}replyLst")
            for rep in list(rl.findall(f"{{{P188}}}reply")) if rl is not None else []:
                if rep.get("id") == comment_id:
                    rl.remove(rep)
            if rl is not None and len(rl) == 0:
                cm.remove(rl)
        ps.set_xml(cpart, cmroot)

    patch_parts(path, mut)
