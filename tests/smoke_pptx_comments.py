"""Smoke test: modern threaded comments in PowerPoint (.pptx).

add -> reply -> resolve/reopen -> delete via _pptx_comments, asserting the authors part, the
per-slide modern comment part, the nested reply (replyLst), the resolved status, byte-stability
of unrelated parts, and that python-pptx still opens each step.

Run: python tests/smoke_pptx_comments.py
"""
import os
import sys
import tempfile
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "scripts"))
import _pptx_comments as pc   # noqa: E402
import _errors                # noqa: E402
from pptx import Presentation   # noqa: E402


def _base(path):
    prs = Presentation()
    for title in ("Quarterly Review Summary", "Regional Breakdown"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])   # Title Only
        slide.shapes.title.text = title
    prs.save(path)


def _opens(path):
    try:
        Presentation(path)
        return True
    except Exception as e:
        return f"{type(e).__name__}: {e}"


def _parts(path):
    with zipfile.ZipFile(path) as z:
        return {i.filename: z.read(i.filename) for i in z.infolist()}


def main() -> int:
    fails = []
    tmp = tempfile.mkdtemp(prefix="pptx_comments_")
    f = os.path.join(tmp, "t.pptx")
    _base(f)
    before = _parts(f)

    try:
        pc.add_comment(f, {"slide": 99}, "x", author="A")
        fails.append("bad slide: expected AnchorNotFound")
    except _errors.AnchorNotFound:
        pass
    except Exception as e:
        fails.append(f"bad slide: wrong error {type(e).__name__}")

    # --- add ---
    cid = pc.add_comment(f, {"slide": 1}, "Please confirm this figure matches the latest forecast.",
                         author="Alex Morgan")
    recs = pc.list_comments(f)
    if len(recs) != 1:
        fails.append(f"add: expected 1 comment, got {len(recs)}")
    else:
        r = recs[0]
        if r["anchor"]["slide"] != 1 or r["location"] != "Slide 1":
            fails.append(f"add: bad anchor/location {r['anchor']} {r['location']}")
        if r["author"] != "Alex Morgan" or r["resolved"] or r["is_reply"]:
            fails.append(f"add: bad record {r}")
        if "Quarterly Review Summary" not in r["context"]:
            fails.append(f"add: context missing slide title: {r['context']!r}")
    with zipfile.ZipFile(f) as z:
        names = set(z.namelist())
        if "ppt/authors.xml" not in names:
            fails.append("add: authors.xml missing")
        if "ppt/comments/modernComment1.xml" not in names:
            fails.append(f"add: modern comment part missing ({[n for n in names if 'comment' in n.lower()]})")
        ct = z.read("[Content_Types].xml").decode()
        for frag in ("ms-powerpoint.comments", "ms-powerpoint.authors"):
            if frag not in ct:
                fails.append(f"add: content-type missing {frag}")
        prels = z.read("ppt/_rels/presentation.xml.rels").decode()
        if "relationships/authors" not in prels:
            fails.append("add: authors relationship missing on presentation")
        srels = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
        if "relationships/comments" not in srels:
            fails.append("add: comments relationship missing on slide")
    if (o := _opens(f)) is not True:
        fails.append(f"add: won't open: {o}")
    changed = sorted(n for n in before if before[n] != _parts(f).get(n))
    allowed = {"[Content_Types].xml", "ppt/_rels/presentation.xml.rels",
               "ppt/slides/slide1.xml", "ppt/slides/_rels/slide1.xml.rels"}
    if not set(changed).issubset(allowed):
        fails.append(f"add: changed unexpected existing parts: {set(changed) - allowed}")

    # --- reply (nested replyLst) ---
    rid = pc.reply(f, cid, "Confirmed - it matches the final forecast.", author="Sam Lee")
    recs = {r["id"]: r for r in pc.list_comments(f)}
    if len(recs) != 2:
        fails.append(f"reply: expected 2, got {len(recs)}")
    if rid in recs:
        if recs[rid]["parent_id"] != cid or recs[rid]["thread_id"] != cid or not recs[rid]["is_reply"]:
            fails.append(f"reply: bad threading {recs[rid]}")
        if recs[rid]["author"] != "Sam Lee":
            fails.append("reply: author wrong")
    else:
        fails.append("reply: reply id not listed")
    with zipfile.ZipFile(f) as z:
        if b"replyLst" not in z.read("ppt/comments/modernComment1.xml"):
            fails.append("reply: no replyLst in comment part")
    if (o := _opens(f)) is not True:
        fails.append(f"reply: won't open: {o}")

    # --- resolve / reopen ---
    pc.set_status(f, cid, True)
    if not all(r["resolved"] for r in pc.list_comments(f)):
        fails.append("resolve: not marked resolved")
    pc.set_status(f, cid, False)
    if any(r["resolved"] for r in pc.list_comments(f)):
        fails.append("reopen: still resolved")

    # --- delete reply then root ---
    pc.delete(f, rid)
    if [r["id"] for r in pc.list_comments(f)] != [cid]:
        fails.append("delete reply: wrong remaining set")
    pc.delete(f, cid)
    if pc.list_comments(f):
        fails.append("delete root: comments remain")
    if (o := _opens(f)) is not True:
        fails.append(f"final: won't open: {o}")

    if fails:
        print("FAIL:")
        for x in fails:
            print("  -", x)
        return 1
    print("PASS — pptx threaded comments (add / reply / resolve / reopen / delete)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
